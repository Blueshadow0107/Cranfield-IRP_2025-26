#!/usr/bin/env python3
"""Populate the MSc CFD poster template with the porous-reactor BZ content."""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
import os

ROOT = "/home/cyanidepopcorn/Cranfield/IndividualResearchProject"
TEMPLATE = os.path.join(ROOT, "Presentations", "MScCFD_Poster_Template.pptx")
OUTPUT = os.path.join(ROOT, "Presentations", "MScCFD_Poster.pptx")
FIG_DIR = os.path.join(ROOT, "Analysis", "figures")

prs = Presentation(TEMPLATE)
slide = prs.slides[0]

# --- Clear existing placeholder shapes (background is on master) ---
for shape in list(slide.shapes):
    sp = shape.element
    sp.getparent().remove(sp)

# --- Helpers ---
BLUE = RGBColor(0, 51, 102)
DARK = RGBColor(33, 33, 33)
WHITE = RGBColor(255, 255, 255)

def add_textbox(left, top, width, height, text, font_size=14, bold=False,
                color=DARK, align=PP_ALIGN.LEFT, italic=False, font_name="Arial"):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = align
    run = p.runs[0]
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = font_name
    return box

def add_bullets(left, top, width, height, bullets, font_size=13, color=DARK):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = bullet
        p.level = 0
        p.alignment = PP_ALIGN.LEFT
        run = p.runs[0]
        run.font.size = Pt(font_size)
        run.font.color.rgb = color
        run.font.name = "Arial"
    return box

def add_section_header(left, top, width, text):
    # Slightly taller header bar with blue text
    return add_textbox(left, top, width, Inches(0.28), text,
                       font_size=18, bold=True, color=BLUE)

def add_picture(left, top, width, path, caption=None, cap_size=10):
    if not os.path.exists(path):
        print(f"WARNING: missing figure {path}")
        return None
    from PIL import Image
    with Image.open(path) as im:
        img_w, img_h = im.size
    aspect = img_h / img_w
    height = width * aspect
    pic = slide.shapes.add_picture(path, left, top, width, height)
    if caption:
        cap_top = top + height + Inches(0.02)
        add_textbox(left, cap_top, width, Inches(0.20), caption,
                    font_size=cap_size, color=DARK, align=PP_ALIGN.CENTER)
    return pic

# --- Layout constants (inches) ---
SLIDE_W = 13.3333
SLIDE_H = 7.5
MARGIN = 0.25
GAP = 0.20
COL_W = (SLIDE_W - 2 * MARGIN - 2 * GAP) / 3  # ~4.11 in
TITLE_H = 0.90
SUBTITLE_H = 0.28
SEC_HEAD_H = 0.28
ROW1_TOP = MARGIN + TITLE_H + SUBTITLE_H + 0.10
ROW_H = (SLIDE_H - ROW1_TOP - MARGIN - 0.10) / 2  # ~2.96 in
ROW2_TOP = ROW1_TOP + ROW_H + 0.10

COL1_L = MARGIN
COL2_L = MARGIN + COL_W + GAP
COL3_L = MARGIN + 2 * (COL_W + GAP)

# --- Title block ---
add_textbox(MARGIN, Inches(0.15), Inches(SLIDE_W - 2 * MARGIN), Inches(0.42),
            "From Measured Pore-Scale Operators to Network-Level Prediction in a Structured Belousov–Zhabotinsky Reactor",
            font_size=26, bold=True, color=BLUE, align=PP_ALIGN.CENTER)
add_textbox(MARGIN, Inches(0.60), Inches(SLIDE_W - 2 * MARGIN), Inches(0.20),
            "Srivijayesh Venugopal  ·  Supervisor: Prof. Weisi Guo  ·  Cranfield University, MSc Computational Fluid Dynamics",
            font_size=14, color=DARK, align=PP_ALIGN.CENTER)

# --- Section 1: Motivation & Porous Reactor Picture (Col1 top) ---
add_section_header(COL1_L, ROW1_TOP, Inches(COL_W), "Motivation & Porous Reactor Picture")
add_bullets(COL1_L, ROW1_TOP + SEC_HEAD_H + 0.02, Inches(COL_W), Inches(ROW_H - SEC_HEAD_H - 0.05),
            [
                "• Digital neural networks face the von Neumann bottleneck: memory and logic are physically separated.",
                "• A structured BZ reactor can act as a programmable analog computer: chemistry itself performs the computation.",
                "• Wave pulses travel through pore channels; collisions, annihilation and timing implement logic.",
                "• Goal: replace discrete node-by-node simulation with a graph of measured pore-scale operators."
            ], font_size=13)

# --- Section 2: Model & Methods (Col2 top) ---
add_section_header(COL2_L, ROW1_TOP, Inches(COL_W), "Model & Methods")
add_bullets(COL2_L, ROW1_TOP + SEC_HEAD_H + 0.02, Inches(COL_W * 0.58), Inches(ROW_H - SEC_HEAD_H - 0.05),
            [
                "• Two-variable Oregonator with light suppression φ and anisotropic diffusion tensor D.",
                "• Explicit finite-difference solver: operator splitting + adaptive reaction subcycling.",
                "• Black-box characterisation: measure speed, refractory window and junction truth tables.",
                "• Inputs: calibrated dark-spot flashes (T_FLASH ≈ 3 time units)."
            ], font_size=12)
# Verification figure on right side of methods section
fig_path = os.path.join(FIG_DIR, "pub", "pub_verification.png")
add_picture(COL2_L + Inches(COL_W * 0.62), ROW1_TOP + SEC_HEAD_H + 0.08,
            Inches(COL_W * 0.36), fig_path,
            caption="MMS verification: 2nd-order convergence", cap_size=9)

# --- Section 3: Measured Pore-Scale Operators (Col3 top) ---
add_section_header(COL3_L, ROW1_TOP, Inches(COL_W), "Measured Pore-Scale Operators")
# Three small figures stacked vertically, bullets beside them
fig_h = (ROW_H - SEC_HEAD_H - 0.15) / 3
pic_w = Inches(COL_W * 0.48)
add_picture(COL3_L, ROW1_TOP + SEC_HEAD_H + 0.03, pic_w,
            os.path.join(FIG_DIR, "pub", "pub_channel_transfer_darkspot.png"),
            caption="Wire: v ≈ 6.46 cells/t.u.", cap_size=8)
add_picture(COL3_L + pic_w + Inches(0.08), ROW1_TOP + SEC_HEAD_H + 0.03, pic_w,
            os.path.join(FIG_DIR, "pub", "pub_tjunction_logic_darkspot.png"),
            caption="Inhibition: A AND (NOT B), >200×", cap_size=8)
add_picture(COL3_L, ROW1_TOP + SEC_HEAD_H + fig_h + 0.10, pic_w * 2 + Inches(0.08),
            os.path.join(FIG_DIR, "pub", "pub_anisotropic_routing.png"),
            caption="Anisotropic routing via sqrt(r) eikonal steering", cap_size=8)

# --- Section 4: Graph Simulator (Col1 bottom) ---
add_section_header(COL1_L, ROW2_TOP, Inches(COL_W), "Graph Simulator")
add_bullets(COL1_L, ROW2_TOP + SEC_HEAD_H + 0.02, Inches(COL_W), Inches(ROW_H - SEC_HEAD_H - 1.65),
            [
                "• Event-driven graph built from measured operators.",
                "• Validation: inhibition gate, priority router, shared-control gate.",
                "• Shared channel introduces crosstalk; geometry-aware routing required."
            ], font_size=13)
add_picture(COL1_L, ROW2_TOP + ROW_H - 1.55, Inches(COL_W),
            os.path.join(FIG_DIR, "rd_multi_gate_pde_snapshots.png"),
            caption="Graph-to-PDE validation: priority router / shared-control gate", cap_size=9)

# --- Section 5: Boundary Experiments (Col2 bottom) ---
add_section_header(COL2_L, ROW2_TOP, Inches(COL_W), "Boundary Experiments")
add_bullets(COL2_L, ROW2_TOP + SEC_HEAD_H + 0.02, Inches(COL_W), Inches(ROW_H - SEC_HEAD_H - 1.85),
            [
                "• 3D inhibition gate: 16.8× separation, speed 6.62 cells/t.u.",
                "• Collision XOR leaks when the merging stem stays excitable (free-medium baseline).",
                "• Diode effect: negative result — symmetry dominates over geometry bias."
            ], font_size=13)
pic_left = add_picture(COL2_L, ROW2_TOP + ROW_H - 1.75, Inches(COL_W * 0.49),
                       os.path.join(FIG_DIR, "rd_3d_transfer_logic_snapshots.png"),
                       caption="3D inhibition gate", cap_size=8)
pic_right = add_picture(COL2_L + Inches(COL_W * 0.51), ROW2_TOP + ROW_H - 1.75, Inches(COL_W * 0.49),
                        os.path.join(FIG_DIR, "rd_spot_xor_protocol.png"),
                        caption="Spot XOR leak protocol", cap_size=8)

# --- Section 6: Conclusions & Future Work (Col3 bottom) ---
add_section_header(COL3_L, ROW2_TOP, Inches(COL_W), "Conclusions & Future Work")
add_bullets(COL3_L, ROW2_TOP + SEC_HEAD_H + 0.02, Inches(COL_W), Inches(ROW_H - SEC_HEAD_H - 0.05),
            [
                "• Pore-network abstraction is viable when junction geometry is carried explicitly.",
                "• Measured operators capture enough physics for fast graph-level prediction.",
                "• Next steps:",
                "  – Curated library of junction shapes (T, Y, diode, collision chamber).",
                "  – Volumetric 3D illumination and full 3D training.",
                "  – Experimental calibration against real photosensitive BZ reactors."
            ], font_size=13)

prs.save(OUTPUT)
print(f"Saved {OUTPUT}")
